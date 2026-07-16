"""
Builder agent — the final node. Turns written slides into a .pptx file.

Reads:  state["written_slides"], state["parsed_doc"], state["audience"]
Writes: state["output_path"], state["current_step"]

Uses python-pptx to generate a real editable PowerPoint file.
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from app.agents.state import AgentState

# Where finished .pptx files land
OUTPUT_DIR = Path("outputs")
OUTPUT_DIR.mkdir(exist_ok=True)


def add_title_slide(prs: Presentation, title: str, subtitle: str):
    """Add a clean title slide as the first slide of the deck."""
    slide_layout = prs.slide_layouts[0]  # 0 = title layout
    slide = prs.slides.add_slide(slide_layout)
    
    slide.shapes.title.text = title
    if len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle


def add_content_slide(prs: Presentation, title: str, bullets: list, speaker_notes: str = None):
    """Add a content slide with title, bullet points, and optional speaker notes."""
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    
    # Guard against absurdly long titles that break layout
    if len(title) > 100:
        title = title[:97] + "..."
    slide.shapes.title.text = title
    
    # Fill in bullets — skip empty ones
    body_placeholder = slide.placeholders[1]
    text_frame = body_placeholder.text_frame
    text_frame.clear()
    
    valid_bullets = [b for b in bullets if b and b.strip()]
    if not valid_bullets:
        valid_bullets = ["(No content generated for this slide)"]
    
    for i, bullet in enumerate(valid_bullets):
        # Cap bullet length to prevent overflow
        clean_bullet = bullet.strip()
        if len(clean_bullet) > 200:
            clean_bullet = clean_bullet[:197] + "..."
        
        if i == 0:
            para = text_frame.paragraphs[0]
        else:
            para = text_frame.add_paragraph()
        
        para.text = clean_bullet
        para.font.size = Pt(20)
        para.level = 0
    
    # Add speaker notes if provided
    if speaker_notes and speaker_notes.strip():
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = speaker_notes.strip()


def builder_agent(state: AgentState) -> dict:
    """
    Builder node. Generates the final .pptx file.
    """
    written = state.get("written_slides")
    doc = state.get("parsed_doc")
    audience = state.get("audience", "student")
    
    if not written or not doc:
        return {
            "errors": state.get("errors", []) + ["Builder: missing written_slides or parsed_doc"],
            "current_step": "builder_failed",
        }
    
    print(f"🎨 [Builder] Building .pptx from {len(written)} slides...")
    
    try:
        # Create a fresh presentation
        prs = Presentation()
        
        # Sort slides by index just in case they arrived out of order
        sorted_slides = sorted(written, key=lambda s: s["index"])
        
        # Build a title slide from the document name + audience
        # Strip job_id prefix (8 hex chars + underscore) and clean up
        doc_name = Path(doc.source_file).stem
        # Remove leading job_id if present (matches "abc12345_" pattern)
        import re
        cleaned = re.sub(r"^[a-f0-9]{8}_", "", doc_name)
        pretty_title = cleaned.replace("_", " ").replace("-", " ").title()
        subtitle = f"Tailored for {audience}s • {len(sorted_slides)} slides"
        add_title_slide(prs, pretty_title, subtitle)
        
        # Build each content slide
        for slide in sorted_slides:
            add_content_slide(
                prs,
                title=slide["title"],
                bullets=slide.get("bullets", []),
                speaker_notes=slide.get("speaker_notes"),
            )
        
        # Save to disk with a descriptive name
        output_name = f"{doc_name}_{audience}_{len(sorted_slides)}slides.pptx"
        output_path = OUTPUT_DIR / output_name
        prs.save(str(output_path))
        
        print(f"✅ [Builder] Saved to: {output_path}")
        
        return {
            "output_path": str(output_path),
            "current_step": "builder_complete",
        }
    
    except Exception as e:
        error_msg = f"Builder failed: {type(e).__name__}: {e}"
        print(f"❌ [Builder] {error_msg}")
        return {
            "output_path": None,
            "errors": state.get("errors", []) + [error_msg],
            "current_step": "builder_failed",
        }